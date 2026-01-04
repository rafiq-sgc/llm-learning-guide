#!/usr/bin/env python3
"""
Create a professional PowerPoint presentation with 15 slides
about "Behind the Scenes of LLMs: How AI Turns English into SQL"
"""

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    # Will be handled in main
    pass

def create_presentation():
    """Create a professional PowerPoint presentation"""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    primary_color = RGBColor(33, 150, 243)  # #2196F3
    secondary_color = RGBColor(76, 175, 80)  # #4CAF50
    accent_color = RGBColor(255, 235, 59)  # #FFEB3B
    danger_color = RGBColor(244, 67, 54)  # #F44336
    dark_bg = RGBColor(26, 26, 26)  # #1a1a1a
    white = RGBColor(255, 255, 255)
    light_gray = RGBColor(245, 245, 245)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    bg = slide1.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(33, 150, 243)
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "🧠 Behind the Scenes of LLMs"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = white
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "How AI Turns English into SQL"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = white
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Meta info
    meta_box = slide1.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.8))
    meta_frame = meta_box.text_frame
    meta_frame.text = "Project: NL2SQL | Duration: 10-15 Minutes"
    meta_para = meta_frame.paragraphs[0]
    meta_para.font.size = Pt(18)
    meta_para.font.color.rgb = white
    meta_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: The Problem
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "The Problem We're Solving"
    title2.text_frame.paragraphs[0].font.size = Pt(44)
    title2.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content2 = slide2.placeholders[1]
    tf2 = content2.text_frame
    tf2.text = "The Challenge"
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(28)
    p2.font.bold = True
    
    # Before/After comparison
    p2 = tf2.add_paragraph()
    p2.text = "Before LLMs:"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.level = 0
    
    p2 = tf2.add_paragraph()
    p2.text = 'User: "Show me all students enrolled in 2024"'
    p2.font.size = Pt(16)
    p2.level = 1
    p2.font.name = "Courier New"
    
    p2 = tf2.add_paragraph()
    p2.text = "Developer: Writes SQL manually"
    p2.font.size = Pt(16)
    p2.level = 1
    
    p2 = tf2.add_paragraph()
    p2.text = "With LLMs: AI automatically generates SQL"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.level = 0
    p2.space_after = Pt(12)
    
    p2 = tf2.add_paragraph()
    p2.text = "Impact:"
    p2.font.size = Pt(20)
    p2.font.bold = True
    
    for impact in ["• Non-technical users can query databases", 
                   "• Faster development", 
                   "• Natural language interface"]:
        p2 = tf2.add_paragraph()
        p2.text = impact
        p2.font.size = Pt(18)
        p2.level = 1
    
    # Slide 3: What is an LLM?
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "What is an LLM?"
    title3.text_frame.paragraphs[0].font.size = Pt(44)
    title3.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = "Large Language Model (LLM)"
    p3 = tf3.paragraphs[0]
    p3.font.size = Pt(28)
    p3.font.bold = True
    
    p3 = tf3.add_paragraph()
    p3.text = "Definition: A neural network trained on massive text data"
    p3.font.size = Pt(18)
    p3.space_after = Pt(12)
    
    p3 = tf3.add_paragraph()
    p3.text = "Key Numbers:"
    p3.font.size = Pt(22)
    p3.font.bold = True
    
    for stat in ["• GPT-3: 175 billion parameters",
                 "• GPT-4: ~1.7 trillion parameters",
                 "• Training Data: 500+ billion tokens"]:
        p3 = tf3.add_paragraph()
        p3.text = stat
        p3.font.size = Pt(18)
        p3.level = 1
    
    p3 = tf3.add_paragraph()
    p3.text = "Analogy: Super Librarian"
    p3.font.size = Pt(22)
    p3.font.bold = True
    p3.space_before = Pt(12)
    
    # Slide 4: Transformer Architecture
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "The Transformer Architecture"
    title4.text_frame.paragraphs[0].font.size = Pt(44)
    title4.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.text = "How LLMs Process Text"
    p4 = tf4.paragraphs[0]
    p4.font.size = Pt(28)
    p4.font.bold = True
    
    # Create a simple flow diagram using shapes
    left = Inches(1)
    top = Inches(3)
    width = Inches(1.5)
    height = Inches(0.8)
    
    shapes = ["Input Text", "Tokenization", "Embedding", "Layers", "Output"]
    for i, shape_text in enumerate(shapes):
        shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                       left + i * Inches(1.6), top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = primary_color if i < 4 else secondary_color
        shape.line.color.rgb = white
        text_frame = shape.text_frame
        text_frame.text = shape_text
        text_frame.paragraphs[0].font.size = Pt(12)
        text_frame.paragraphs[0].font.color.rgb = white
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Add arrow
        if i < len(shapes) - 1:
            arrow = slide4.shapes.add_connector(1, left + i * Inches(1.6) + width, 
                                               top + height/2, 
                                               left + (i+1) * Inches(1.6), 
                                               top + height/2)
            arrow.line.color.rgb = dark_bg
            arrow.line.width = Pt(2)
    
    p4 = tf4.add_paragraph()
    p4.text = "Core Components:"
    p4.font.size = Pt(20)
    p4.font.bold = True
    p4.space_before = Pt(24)
    
    for comp in ["1. Tokenization: Convert text to numbers",
                 "2. Embeddings: Convert tokens to vectors",
                 "3. Attention: Focus on relevant parts",
                 "4. Layers: 96+ layers of processing"]:
        p4 = tf4.add_paragraph()
        p4.text = comp
        p4.font.size = Pt(16)
        p4.level = 1
    
    # Slide 5: Tokenization
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "Tokenization"
    title5.text_frame.paragraphs[0].font.size = Pt(44)
    title5.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content5 = slide5.placeholders[1]
    tf5 = content5.text_frame
    tf5.text = "Converting Words to Numbers"
    p5 = tf5.paragraphs[0]
    p5.font.size = Pt(28)
    p5.font.bold = True
    
    p5 = tf5.add_paragraph()
    p5.text = "We see: \"Show me students\""
    p5.font.size = Pt(18)
    p5.font.name = "Courier New"
    p5.space_after = Pt(12)
    
    p5 = tf5.add_paragraph()
    p5.text = "LLM sees: [1234, 567, 8901]"
    p5.font.size = Pt(18)
    p5.font.name = "Courier New"
    p5.space_after = Pt(12)
    
    p5 = tf5.add_paragraph()
    p5.text = "Process:"
    p5.font.size = Pt(20)
    p5.font.bold = True
    
    for step in ["1. Text → Split into tokens",
                 "2. Tokens → Convert to numbers",
                 "3. Numbers → Process through neural network"]:
        p5 = tf5.add_paragraph()
        p5.text = step
        p5.font.size = Pt(18)
        p5.level = 1
    
    # Slide 6: Attention Mechanism
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    title6.text = "Attention Mechanism"
    title6.text_frame.paragraphs[0].font.size = Pt(44)
    title6.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content6 = slide6.placeholders[1]
    tf6 = content6.text_frame
    tf6.text = "The Magic"
    p6 = tf6.paragraphs[0]
    p6.font.size = Pt(28)
    p6.font.bold = True
    
    p6 = tf6.add_paragraph()
    p6.text = "The Problem: \"The cat sat on the mat because it was tired\""
    p6.font.size = Pt(16)
    p6.space_after = Pt(12)
    
    p6 = tf6.add_paragraph()
    p6.text = "What does 'it' refer to?"
    p6.font.size = Pt(18)
    p6.font.bold = True
    p6.space_after = Pt(12)
    
    p6 = tf6.add_paragraph()
    p6.text = "The Solution: Attention"
    p6.font.size = Pt(20)
    p6.font.bold = True
    
    p6 = tf6.add_paragraph()
    p6.text = "Attention(\"it\" to \"cat\") = 0.9 ← High!"
    p6.font.size = Pt(16)
    p6.font.name = "Courier New"
    p6.level = 1
    
    p6 = tf6.add_paragraph()
    p6.text = "Multi-Head Attention:"
    p6.font.size = Pt(20)
    p6.font.bold = True
    p6.space_before = Pt(12)
    
    for head in ["• Head 1: Grammar (subject-verb)",
                 "• Head 2: Meaning (cat → animal)",
                 "• Head 3: Long-range dependencies"]:
        p6 = tf6.add_paragraph()
        p6.text = head
        p6.font.size = Pt(18)
        p6.level = 1
    
    # Slide 7: Text Generation
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    title7.text = "How LLMs Generate Text"
    title7.text_frame.paragraphs[0].font.size = Pt(44)
    title7.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content7 = slide7.placeholders[1]
    tf7 = content7.text_frame
    tf7.text = "The Generation Process"
    p7 = tf7.paragraphs[0]
    p7.font.size = Pt(28)
    p7.font.bold = True
    
    p7 = tf7.add_paragraph()
    p7.text = "Step-by-Step:"
    p7.font.size = Pt(20)
    p7.font.bold = True
    
    steps = ["1. Input Processing: Tokenize and embed",
            "2. Layer Processing: 96 transformer layers",
            "3. Next Token Prediction: Probability distribution",
            "4. Sampling: Choose next token",
            "5. Iterate: Repeat until complete"]
    
    for step in steps:
        p7 = tf7.add_paragraph()
        p7.text = step
        p7.font.size = Pt(18)
        p7.level = 1
    
    p7 = tf7.add_paragraph()
    p7.text = "Example Probability:"
    p7.font.size = Pt(20)
    p7.font.bold = True
    p7.space_before = Pt(12)
    
    p7 = tf7.add_paragraph()
    p7.text = '"enrolled" = 0.35 ← Most likely'
    p7.font.size = Pt(16)
    p7.font.name = "Courier New"
    p7.level = 1
    
    # Slide 8: Training
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    title8.text = "Training LLMs"
    title8.text_frame.paragraphs[0].font.size = Pt(44)
    title8.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content8 = slide8.placeholders[1]
    tf8 = content8.text_frame
    tf8.text = "Two Stages"
    p8 = tf8.paragraphs[0]
    p8.font.size = Pt(28)
    p8.font.bold = True
    
    # Create training flow
    left = Inches(1)
    top = Inches(3.5)
    width = Inches(2)
    height = Inches(0.8)
    
    stages = ["Massive Data", "Pre-training", "Base LLM", "Fine-tuning", "Specialized"]
    colors = [light_gray, accent_color, primary_color, secondary_color, secondary_color]
    
    for i, (stage, color) in enumerate(zip(stages, colors)):
        shape = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       left + i * Inches(1.6), top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = dark_bg
        text_frame = shape.text_frame
        text_frame.text = stage
        text_frame.paragraphs[0].font.size = Pt(11)
        text_frame.paragraphs[0].font.color.rgb = dark_bg if i == 0 else white
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        if i < len(stages) - 1:
            arrow = slide8.shapes.add_connector(1, left + i * Inches(1.6) + width,
                                               top + height/2,
                                               left + (i+1) * Inches(1.6),
                                               top + height/2)
            arrow.line.color.rgb = dark_bg
            arrow.line.width = Pt(2)
    
    p8 = tf8.add_paragraph()
    p8.text = "Stage 1: Pre-training (500B+ tokens, millions of dollars)"
    p8.font.size = Pt(18)
    p8.space_before = Pt(24)
    
    p8 = tf8.add_paragraph()
    p8.text = "Stage 2: Fine-tuning (Task-specific examples)"
    p8.font.size = Pt(18)
    
    # Slide 9: Why LLMs Succeed
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    title9.text = "Why LLMs Succeed"
    title9.text_frame.paragraphs[0].font.size = Pt(44)
    title9.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content9 = slide9.placeholders[1]
    tf9 = content9.text_frame
    tf9.text = "The Success Factors"
    p9 = tf9.paragraphs[0]
    p9.font.size = Pt(28)
    p9.font.bold = True
    
    factors = [
        "1. Scale Matters: GPT-1 (117M) → GPT-4 (~1.7T parameters)",
        "2. Massive Training Data: Entire internet, books, code",
        "3. Architecture Advantages: Parallelization, attention, transfer learning",
        "4. In-Context Learning: Learns from examples in prompt"
    ]
    
    for factor in factors:
        p9 = tf9.add_paragraph()
        p9.text = factor
        p9.font.size = Pt(18)
        p9.level = 0
        p9.space_after = Pt(8)
    
    # Slide 10: Why LLMs Fail
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    title10.text = "Why LLMs Fail"
    title10.text_frame.paragraphs[0].font.size = Pt(44)
    title10.text_frame.paragraphs[0].font.color.rgb = danger_color
    
    content10 = slide10.placeholders[1]
    tf10 = content10.text_frame
    tf10.text = "Common Failure Modes"
    p10 = tf10.paragraphs[0]
    p10.font.size = Pt(28)
    p10.font.bold = True
    
    failures = [
        "1. Hallucination: Makes up plausible but incorrect information",
        "2. Context Limits: GPT-4 has 8,192 token limit",
        "3. Ambiguity: Multiple interpretations possible",
        "4. Schema Mismatch: Assumes schema that doesn't exist"
    ]
    
    for failure in failures:
        p10 = tf10.add_paragraph()
        p10.text = failure
        p10.font.size = Pt(18)
        p10.level = 0
        p10.space_after = Pt(8)
    
    # Slide 11: NL2SQL Challenge
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title11 = slide11.shapes.title
    title11.text = "NL2SQL - The Challenge"
    title11.text_frame.paragraphs[0].font.size = Pt(44)
    title11.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content11 = slide11.placeholders[1]
    tf11 = content11.text_frame
    tf11.text = "Converting Natural Language to SQL"
    p11 = tf11.paragraphs[0]
    p11.font.size = Pt(28)
    p11.font.bold = True
    
    p11 = tf11.add_paragraph()
    p11.text = 'Input: "Show me all students enrolled in 2024"'
    p11.font.size = Pt(18)
    p11.font.name = "Courier New"
    p11.space_after = Pt(12)
    
    p11 = tf11.add_paragraph()
    p11.text = "Output: SELECT * FROM students WHERE enrollment_year = 2024"
    p11.font.size = Pt(16)
    p11.font.name = "Courier New"
    p11.space_after = Pt(12)
    
    p11 = tf11.add_paragraph()
    p11.text = "Why It's Hard:"
    p11.font.size = Pt(20)
    p11.font.bold = True
    
    for reason in ["• Language Ambiguity",
                   "• Schema Complexity",
                   "• Domain Knowledge Required",
                   "• SQL Expertise Needed"]:
        p11 = tf11.add_paragraph()
        p11.text = reason
        p11.font.size = Pt(18)
        p11.level = 1
    
    # Slide 12: System Architecture
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    title12 = slide12.shapes.title
    title12.text = "Your NL2SQL System Architecture"
    title12.text_frame.paragraphs[0].font.size = Pt(40)
    title12.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content12 = slide12.placeholders[1]
    tf12 = content12.text_frame
    tf12.text = "Complete Flow"
    p12 = tf12.paragraphs[0]
    p12.font.size = Pt(28)
    p12.font.bold = True
    
    # Create architecture flow
    left = Inches(0.5)
    top = Inches(3.5)
    width = Inches(1.2)
    height = Inches(0.6)
    
    arch_components = ["User", "Intent", "Hybrid", "Context", "Agent", "LLM", "SQL", "Results"]
    arch_colors = [primary_color, secondary_color, accent_color, primary_color, 
                   secondary_color, primary_color, secondary_color, secondary_color]
    
    for i, (comp, color) in enumerate(zip(arch_components, arch_colors)):
        shape = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        left + i * Inches(1.1), top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = white
        text_frame = shape.text_frame
        text_frame.text = comp
        text_frame.paragraphs[0].font.size = Pt(10)
        text_frame.paragraphs[0].font.color.rgb = white
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        if i < len(arch_components) - 1:
            arrow = slide12.shapes.add_connector(1, left + i * Inches(1.1) + width,
                                                top + height/2,
                                                left + (i+1) * Inches(1.1),
                                                top + height/2)
            arrow.line.color.rgb = dark_bg
            arrow.line.width = Pt(2)
    
    p12 = tf12.add_paragraph()
    p12.text = "Key Components:"
    p12.font.size = Pt(20)
    p12.font.bold = True
    p12.space_before = Pt(24)
    
    for comp in ["1. Intent Classification",
                 "2. Hybrid Search (Vector + BM25)",
                 "3. SQL Agent (ReAct Pattern)",
                 "4. LLM + Tools"]:
        p12 = tf12.add_paragraph()
        p12.text = comp
        p12.font.size = Pt(16)
        p12.level = 1
    
    # Slide 13: Hybrid Search
    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    title13 = slide13.shapes.title
    title13.text = "Hybrid Search"
    title13.text_frame.paragraphs[0].font.size = Pt(44)
    title13.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content13 = slide13.placeholders[1]
    tf13 = content13.text_frame
    tf13.text = "Finding Relevant Tables"
    p13 = tf13.paragraphs[0]
    p13.font.size = Pt(28)
    p13.font.bold = True
    
    # Create hybrid search diagram
    left = Inches(2)
    top = Inches(3)
    
    # User Query
    query_shape = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          left + Inches(2), top, Inches(2), Inches(0.6))
    query_shape.fill.solid()
    query_shape.fill.fore_color.rgb = primary_color
    query_shape.text_frame.text = "User Query"
    query_shape.text_frame.paragraphs[0].font.size = Pt(12)
    query_shape.text_frame.paragraphs[0].font.color.rgb = white
    query_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Vector Search
    vector_shape = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            left, top + Inches(1.2), Inches(2), Inches(0.6))
    vector_shape.fill.solid()
    vector_shape.fill.fore_color.rgb = secondary_color
    vector_shape.text_frame.text = "Vector Search"
    vector_shape.text_frame.paragraphs[0].font.size = Pt(12)
    vector_shape.text_frame.paragraphs[0].font.color.rgb = white
    vector_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # BM25 Search
    bm25_shape = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         left + Inches(4), top + Inches(1.2), Inches(2), Inches(0.6))
    bm25_shape.fill.solid()
    bm25_shape.fill.fore_color.rgb = accent_color
    bm25_shape.text_frame.text = "BM25 Search"
    bm25_shape.text_frame.paragraphs[0].font.size = Pt(12)
    bm25_shape.text_frame.paragraphs[0].font.color.rgb = dark_bg
    bm25_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Ensemble
    ensemble_shape = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              left + Inches(2), top + Inches(2.4), Inches(2), Inches(0.6))
    ensemble_shape.fill.solid()
    ensemble_shape.fill.fore_color.rgb = primary_color
    ensemble_shape.text_frame.text = "Ensemble"
    ensemble_shape.text_frame.paragraphs[0].font.size = Pt(12)
    ensemble_shape.text_frame.paragraphs[0].font.color.rgb = white
    ensemble_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add connectors
    arrow1 = slide13.shapes.add_connector(1, left + Inches(3), top + Inches(0.6),
                                         left + Inches(1), top + Inches(1.2))
    arrow1.line.color.rgb = dark_bg
    arrow1.line.width = Pt(2)
    
    arrow2 = slide13.shapes.add_connector(1, left + Inches(3), top + Inches(0.6),
                                         left + Inches(5), top + Inches(1.2))
    arrow2.line.color.rgb = dark_bg
    arrow2.line.width = Pt(2)
    
    arrow3 = slide13.shapes.add_connector(1, left + Inches(1), top + Inches(1.8),
                                         left + Inches(3), top + Inches(2.4))
    arrow3.line.color.rgb = dark_bg
    arrow3.line.width = Pt(2)
    
    arrow4 = slide13.shapes.add_connector(1, left + Inches(5), top + Inches(1.8),
                                         left + Inches(3), top + Inches(2.4))
    arrow4.line.color.rgb = dark_bg
    arrow4.line.width = Pt(2)
    
    p13 = tf13.add_paragraph()
    p13.text = "Why Hybrid?"
    p13.font.size = Pt(20)
    p13.font.bold = True
    p13.space_before = Pt(24)
    
    p13 = tf13.add_paragraph()
    p13.text = "Vector: Handles synonyms | BM25: Exact matches | Together: Best of both!"
    p13.font.size = Pt(16)
    
    # Slide 14: ReAct Agent
    slide14 = prs.slides.add_slide(prs.slide_layouts[1])
    title14 = slide14.shapes.title
    title14.text = "SQL Agent - ReAct Pattern"
    title14.text_frame.paragraphs[0].font.size = Pt(40)
    title14.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content14 = slide14.placeholders[1]
    tf14 = content14.text_frame
    tf14.text = "Reasoning + Acting"
    p14 = tf14.paragraphs[0]
    p14.font.size = Pt(28)
    p14.font.bold = True
    
    p14 = tf14.add_paragraph()
    p14.text = "ReAct = Reasoning + Acting"
    p14.font.size = Pt(20)
    p14.font.bold = True
    p14.space_after = Pt(12)
    
    react_steps = [
        "1. User asks question",
        "2. Agent thinks: 'I need to find tables'",
        "3. Agent acts: Uses tool to list tables",
        "4. Agent observes: Gets table list",
        "5. Agent thinks: 'Now check schema'",
        "6. Agent acts: Gets schema",
        "7. Agent generates SQL",
        "8. Agent executes and returns results"
    ]
    
    for step in react_steps:
        p14 = tf14.add_paragraph()
        p14.text = step
        p14.font.size = Pt(16)
        p14.level = 1
    
    # Slide 15: Complete Example & Takeaways
    slide15 = prs.slides.add_slide(prs.slide_layouts[1])
    title15 = slide15.shapes.title
    title15.text = "Complete Example & Key Takeaways"
    title15.text_frame.paragraphs[0].font.size = Pt(36)
    title15.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content15 = slide15.placeholders[1]
    tf15 = content15.text_frame
    tf15.text = "Example Walkthrough"
    p15 = tf15.paragraphs[0]
    p15.font.size = Pt(28)
    p15.font.bold = True
    
    p15 = tf15.add_paragraph()
    p15.text = 'User: "Show me all students enrolled in 2024"'
    p15.font.size = Pt(18)
    p15.font.name = "Courier New"
    p15.space_after = Pt(12)
    
    example_steps = [
        "1. Intent: SQL (0.99 confidence)",
        "2. Hybrid Search: Finds 'students' table",
        "3. SQL Agent: Generates SQL",
        "4. Execute: Returns results",
        "5. Format: Structured response"
    ]
    
    for step in example_steps:
        p15 = tf15.add_paragraph()
        p15.text = step
        p15.font.size = Pt(18)
        p15.level = 1
    
    p15 = tf15.add_paragraph()
    p15.text = "Key Takeaways:"
    p15.font.size = Pt(22)
    p15.font.bold = True
    p15.space_before = Pt(12)
    
    takeaways = [
        "• LLMs are sophisticated pattern matching systems",
        "• Scale matters: More parameters = better performance",
        "• Context is crucial for good results",
        "• Your system combines multiple techniques",
        "• Validation prevents errors"
    ]
    
    for takeaway in takeaways:
        p15 = tf15.add_paragraph()
        p15.text = takeaway
        p15.font.size = Pt(18)
        p15.level = 1
    
    # Save presentation
    output_file = "LLM_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ Presentation created successfully: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"📁 File location: {os.path.abspath(output_file)}")
    
    return output_file

if __name__ == "__main__":
    # Check if python-pptx is installed
    try:
        import pptx
        print("✅ python-pptx is installed")
    except ImportError:
        print("❌ Error: python-pptx not installed")
        print("📦 Installing python-pptx...")
        import subprocess
        import sys
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "--quiet"])
            print("✅ python-pptx installed successfully")
        except Exception as install_error:
            print(f"❌ Failed to install python-pptx: {install_error}")
            print("📦 Please install manually: pip install python-pptx")
            sys.exit(1)
    
    # Create presentation
    try:
        output_file = create_presentation()
        print(f"\n🎉 Success! Presentation created: {output_file}")
        print(f"📂 Location: {os.path.abspath(output_file)}")
        print(f"📊 Total slides: 15")
        print("\n💡 Next steps:")
        print("   1. Open the .pptx file in PowerPoint")
        print("   2. Review and customize as needed")
        print("   3. Add your name and date to slide 1")
        print("   4. Practice your presentation!")
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)

